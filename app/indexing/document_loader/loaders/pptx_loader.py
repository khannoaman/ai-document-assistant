from langchain_core.documents import Document
from pptx import Presentation


def load(file_path: str) -> list[Document]:
    """No lightweight LangChain-native PPTX loader exists (only
    UnstructuredPowerPointLoader, which pulls in numba/llvmlite). This wraps
    python-pptx directly and returns standard Document objects instead.
    """
    presentation = Presentation(file_path)
    documents = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        documents.append(
            Document(
                page_content="\n".join(text_parts),
                metadata={
                    "source": file_path,
                    "slide_number": slide_number,
                    "loader_used": "python-pptx",
                },
            )
        )
    return documents
